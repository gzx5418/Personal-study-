from __future__ import annotations

from io import BytesIO
from pathlib import Path
import html
import os
import re
import shutil
import subprocess
import sys
import time
from textwrap import wrap

from config import settings


SLIDE_W = 1280
SLIDE_H = 720


def _plain_text(text: str) -> str:
    text = re.sub(r"`([^`]+)`", r"\1", text)
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    text = re.sub(r"[*_>#-]+", "", text)
    return text.strip()


def _safe_filename(value: str) -> str:
    safe = re.sub(r"[^0-9A-Za-z\u4e00-\u9fff_-]+", "_", value.strip())
    safe = re.sub(r"_+", "_", safe).strip("_")
    return safe[:80] or "learning_deck"


def _split_text(text: str, max_chars: int) -> list[str]:
    text = _plain_text(text)
    if not text:
        return []
    # textwrap does not understand CJK widths, but handles long English tokens well enough.
    lines: list[str] = []
    for paragraph in re.split(r"\s*\n\s*", text):
        if not paragraph:
            continue
        if len(paragraph) <= max_chars:
            lines.append(paragraph)
        else:
            lines.extend(wrap(paragraph, width=max_chars, break_long_words=False) or [paragraph[:max_chars]])
    return lines


def _parse_outline(outline_md: str, title: str) -> list[dict[str, object]]:
    sections: list[dict[str, object]] = []
    current_title = ""
    current_lines: list[str] = []

    for raw_line in outline_md.splitlines():
        line = raw_line.strip()
        if line.startswith("## "):
            if current_title:
                sections.append({"title": current_title, "bullets": current_lines})
            current_title = _plain_text(line[3:])
            current_lines = []
        elif current_title and line:
            cleaned = _plain_text(line)
            if cleaned:
                current_lines.append(cleaned)

    if current_title:
        sections.append({"title": current_title, "bullets": current_lines})

    if not sections:
        sections = [{"title": title or "个性化学习资源", "bullets": _split_text(outline_md, 42)[:8]}]

    return sections[:18]


def _svg_text(x: int, y: int, text: str, size: int, weight: str = "400", color: str = "#1f2937") -> str:
    return (
        f'<text x="{x}" y="{y}" font-family="Microsoft YaHei, SimHei, Arial" '
        f'font-size="{size}" font-weight="{weight}" fill="{color}">{html.escape(text)}</text>'
    )


def _render_cover_svg(title: str, subtitle: str) -> str:
    title_lines = _split_text(title or "个性化学习资源", 18)[:2]
    subtitle_lines = _split_text(subtitle, 32)[:2]
    title_svg = "\n".join(_svg_text(90, 255 + i * 70, line, 52, "700", "#111827") for i, line in enumerate(title_lines))
    subtitle_svg = "\n".join(_svg_text(94, 430 + i * 34, line, 24, "400", "#475569") for i, line in enumerate(subtitle_lines))
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{SLIDE_W}" height="{SLIDE_H}" viewBox="0 0 {SLIDE_W} {SLIDE_H}">
  <rect width="{SLIDE_W}" height="{SLIDE_H}" fill="#f8fafc"/>
  <g id="background">
    <rect x="0" y="0" width="{SLIDE_W}" height="{SLIDE_H}" fill="#f8fafc"/>
    <path d="M860 0 H1280 V720 H680 C820 565 890 410 860 0Z" fill="#dbeafe"/>
    <circle cx="1050" cy="180" r="170" fill="#fef3c7"/>
  </g>
  <g id="header">
    <rect x="90" y="82" width="156" height="38" rx="19" fill="#2563eb"/>
    {_svg_text(112, 108, "智学助手", 20, "700", "#ffffff")}
  </g>
  <g id="title">
    {title_svg}
    {subtitle_svg}
  </g>
  <g id="footer">
    {_svg_text(94, 650, "基于学习画像与课程知识库自动生成", 18, "400", "#64748b")}
  </g>
</svg>'''


def _render_content_svg(page_title: str, bullets: list[str], page_num: int) -> str:
    bullet_lines = []
    y = 180
    for idx, bullet in enumerate(bullets[:7]):
        lines = _split_text(bullet, 44)[:2]
        if not lines:
            continue
        bullet_lines.append(f'<circle cx="112" cy="{y - 8}" r="7" fill="#2563eb"/>')
        bullet_lines.append(_svg_text(136, y, lines[0], 25, "500", "#1f2937"))
        if len(lines) > 1:
            bullet_lines.append(_svg_text(136, y + 33, lines[1], 21, "400", "#475569"))
            y += 76
        else:
            y += 54
    content = "\n    ".join(bullet_lines)
    return f'''<svg xmlns="http://www.w3.org/2000/svg" width="{SLIDE_W}" height="{SLIDE_H}" viewBox="0 0 {SLIDE_W} {SLIDE_H}">
  <rect width="{SLIDE_W}" height="{SLIDE_H}" fill="#ffffff"/>
  <g id="background">
    <rect x="0" y="0" width="{SLIDE_W}" height="{SLIDE_H}" fill="#ffffff"/>
    <rect x="0" y="0" width="32" height="{SLIDE_H}" fill="#2563eb"/>
    <rect x="90" y="610" width="1080" height="1" fill="#e2e8f0"/>
  </g>
  <g id="title">
    {_svg_text(90, 105, page_title[:42], 36, "700", "#0f172a")}
    <rect x="90" y="128" width="78" height="6" rx="3" fill="#f59e0b"/>
  </g>
  <g id="body">
    {content}
  </g>
  <g id="footer">
    {_svg_text(90, 655, "智学助手 · 个性化学习资源", 16, "400", "#64748b")}
    {_svg_text(1140, 655, str(page_num), 18, "700", "#2563eb")}
  </g>
</svg>'''


def _write_ppt_master_project(project_path: Path, outline_md: str, title: str) -> None:
    for dirname in ("svg_output", "svg_final", "notes", "exports", "sources", "images", "templates"):
        (project_path / dirname).mkdir(parents=True, exist_ok=True)

    sections = _parse_outline(outline_md, title)
    source_path = project_path / "sources" / "outline.md"
    source_path.write_text(outline_md, encoding="utf-8")

    design_spec = (
        "# Design Specification\n\n"
        "- Format: ppt169\n"
        "- Style: clean academic education deck\n"
        "- Colors: blue, amber, slate, white\n"
        "- Typography: Microsoft YaHei / SimHei fallback\n"
        "- Generated by backend/services/pptx_service.py using ppt-master SVG export pipeline\n"
    )
    (project_path / "design_spec.md").write_text(design_spec, encoding="utf-8")
    (project_path / "README.md").write_text(f"# {title}\n\n- Canvas format: ppt169\n", encoding="utf-8")

    notes_parts = []
    cover_name = "slide_01_cover"
    (project_path / "svg_output" / f"{cover_name}.svg").write_text(
        _render_cover_svg(title, "面向当前学习画像的课程讲解材料"),
        encoding="utf-8",
    )
    notes_parts.append(f"# {cover_name}\n\n开场说明本节主题、学习目标和使用方式。\n")

    for idx, section in enumerate(sections, start=2):
        stem = f"slide_{idx:02d}"
        page_title = str(section["title"])
        bullets = [str(item) for item in section.get("bullets", [])]
        (project_path / "svg_output" / f"{stem}.svg").write_text(
            _render_content_svg(page_title, bullets, idx),
            encoding="utf-8",
        )
        notes_text = "\n".join(f"- {line}" for line in bullets[:6]) or "围绕本页标题展开讲解。"
        notes_parts.append(f"# {stem}\n\n{notes_text}\n")

    (project_path / "notes" / "total.md").write_text("\n---\n\n".join(notes_parts), encoding="utf-8")


def _run_ppt_master_export(project_path: Path, output_path: Path) -> None:
    ppt_master_dir = Path(settings.PPT_MASTER_DIR)
    script_dir = ppt_master_dir / "skills" / "ppt-master" / "scripts"
    total_split = script_dir / "total_md_split.py"
    export_script = script_dir / "svg_to_pptx.py"

    if not export_script.exists() or not total_split.exists():
        raise FileNotFoundError("ppt-master scripts not found")

    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"
    subprocess.run(
        [sys.executable, str(total_split), str(project_path), "-q"],
        cwd=str(ppt_master_dir),
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        env=env,
    )
    # Native-only export keeps the deck editable and avoids optional PNG fallback renderers.
    subprocess.run(
        [
            sys.executable,
            str(export_script),
            str(project_path),
            "-s",
            "output",
            "--only",
            "native",
            "--no-compat",
            "-a",
            "mixed",
            "--animation-trigger",
            "after-previous",
            "-o",
            str(output_path),
            "-q",
        ],
        cwd=str(ppt_master_dir),
        check=True,
        capture_output=True,
        text=True,
        encoding="utf-8",
        errors="replace",
        env=env,
    )


def generate_pptx_with_ppt_master(outline_md: str, title: str) -> bytes:
    project_root = Path(os.path.dirname(settings.PROFILE_FILE)) / "ppt_master_projects"
    project_root.mkdir(parents=True, exist_ok=True)
    project_name = f"{_safe_filename(title)}_{int(time.time())}"
    project_path = project_root / project_name
    if project_path.exists():
        shutil.rmtree(project_path)
    project_path.mkdir(parents=True, exist_ok=True)

    _write_ppt_master_project(project_path, outline_md, title)
    output_path = project_path / "exports" / f"{_safe_filename(title)}.pptx"
    _run_ppt_master_export(project_path, output_path)
    return output_path.read_bytes()


def generate_pptx_simple(outline_md: str, title: str) -> bytes:
    """Fallback converter when ppt-master scripts or dependencies are unavailable."""
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    prs.core_properties.title = title or "学习资源"

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title or "个性化学习资源"
    title_slide.placeholders[1].text = "由智学助手根据学习画像与课程知识库生成"

    for section in _parse_outline(outline_md, title):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(section["title"])[:80]
        body = slide.placeholders[1].text_frame
        body.clear()
        bullets = [str(item) for item in section.get("bullets", []) if str(item).strip()][:8]
        if not bullets:
            bullets = ["请结合讲者提示进行课堂讲解。"]
        for idx, line in enumerate(bullets):
            paragraph = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            paragraph.text = line[:180]
            paragraph.level = 0
            paragraph.font.size = Pt(18)

    output = BytesIO()
    prs.save(output)
    return output.getvalue()


def generate_pptx_from_outline(outline_md: str, title: str) -> bytes:
    try:
        return generate_pptx_with_ppt_master(outline_md, title)
    except Exception:
        return generate_pptx_simple(outline_md, title)
